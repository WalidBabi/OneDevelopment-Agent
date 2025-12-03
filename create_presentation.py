#!/usr/bin/env python3
"""
PowerPoint Presentation Generator: 20 Agents in 30 Days
Customized for OneDevelopment (OneUAE) - Real Estate Developer
Using DeepAgents Framework, LangSmith Agent Builder & Cursor
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Color palette - OneDevelopment Purple Theme
COLORS = {
    'dark_bg': RGBColor(26, 13, 48),  # Deep purple background
    'accent_purple': RGBColor(150, 107, 252),  # OneDev brand purple
    'accent_gold': RGBColor(255, 193, 7),  # Gold accent
    'accent_cyan': RGBColor(34, 211, 238),
    'accent_green': RGBColor(52, 211, 153),
    'accent_orange': RGBColor(251, 146, 60),
    'accent_pink': RGBColor(244, 114, 182),
    'text_white': RGBColor(255, 255, 255),
    'text_gray': RGBColor(180, 180, 195),
    'card_bg': RGBColor(52, 26, 96),  # OneDev purple #341a60
}

def set_slide_background(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "20 AGENTS IN 30 DAYS"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI-Powered Automation for OneDevelopment"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['accent_gold']
    p.alignment = PP_ALIGN.CENTER
    
    # Tech stack line
    tech_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.6))
    tf = tech_box.text_frame
    p = tf.paragraphs[0]
    p.text = "DeepAgents Framework  •  LangSmith Agent Builder  •  Cursor AI"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['accent_purple']
    p.alignment = PP_ALIGN.CENTER
    
    # Company & Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "OneUAE  |  December 2025"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text_gray']
    p.alignment = PP_ALIGN.CENTER

def add_agenda_slide(prs):
    """Slide 2: Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📋 AGENDA"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    
    agenda_items = [
        ("01", "Why AI Agents for OneDevelopment?", COLORS['accent_purple']),
        ("02", "Our Tech Stack & Tools", COLORS['accent_cyan']),
        ("03", "30-Day Implementation Roadmap", COLORS['accent_green']),
        ("04", "20 Agents for Employee Productivity", COLORS['accent_gold']),
        ("05", "Expected Impact & ROI", COLORS['accent_pink']),
    ]
    
    y_pos = 1.8
    for num, text, color in agenda_items:
        num_box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(0.8), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = color
        
        text_box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(6), Inches(0.6))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['text_white']
        
        y_pos += 0.85

def add_why_agents_slide(prs):
    """Slide 3: Why AI Agents for OneDevelopment"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🏢 WHY AI AGENTS FOR ONEDEVELOPMENT?"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    
    # Big numbers section
    stats = [
        ("20", "AGENTS", COLORS['accent_purple']),
        ("30", "DAYS", COLORS['accent_gold']),
        ("500+", "HRS/MONTH", COLORS['accent_green']),
    ]
    
    x_pos = 0.8
    for num, label, color in stats:
        num_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.4), Inches(2.8), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        label_box = slide.shapes.add_textbox(Inches(x_pos), Inches(2.4), Inches(2.8), Inches(0.4))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_gray']
        p.alignment = PP_ALIGN.CENTER
        
        x_pos += 3
    
    # Key challenges solved
    challenges_title = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(0.4))
    tf = challenges_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Challenges We're Solving:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_cyan']
    
    challenges = [
        "Manual property data entry and market research",
        "Repetitive client inquiries across sales teams",
        "Document processing for contracts and agreements",
        "Project tracking across multiple developments",
        "Employee onboarding and HR processes"
    ]
    
    y_pos = 3.5
    for challenge in challenges:
        point_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.4))
        tf = point_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"✓  {challenge}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS['text_white']
        y_pos += 0.4

def add_tech_stack_slide(prs):
    """Slide 4: Tech Stack Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🔧 OUR DEVELOPMENT STACK"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    
    tools = [
        {
            'name': 'DeepAgents',
            'color': COLORS['accent_purple'],
            'desc': 'Complex Multi-Step Tasks',
            'features': [
                '• Planning & task breakdown',
                '• Persistent memory',
                '• Subagent delegation',
                '• UAE market context'
            ]
        },
        {
            'name': 'LangSmith Builder',
            'color': COLORS['accent_gold'],
            'desc': 'Rapid No-Code Agents',
            'features': [
                '• Natural language config',
                '• Built-in learning',
                '• Pre-built templates',
                '• API integrations'
            ]
        },
        {
            'name': 'Cursor AI',
            'color': COLORS['accent_cyan'],
            'desc': '10x Dev Speed',
            'features': [
                '• AI code generation',
                '• Multi-file editing',
                '• Shell automation',
                '• Real-time assistance'
            ]
        }
    ]
    
    x_positions = [0.3, 3.4, 6.5]
    
    for i, tool in enumerate(tools):
        x = x_positions[i]
        
        name_box = slide.shapes.add_textbox(Inches(x), Inches(1.1), Inches(3), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = tool['name']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = tool['color']
        
        desc_box = slide.shapes.add_textbox(Inches(x), Inches(1.5), Inches(3), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = tool['desc']
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['text_gray']
        
        y_pos = 2
        for feature in tool['features']:
            feat_box = slide.shapes.add_textbox(Inches(x), Inches(y_pos), Inches(3), Inches(0.35))
            tf = feat_box.text_frame
            p = tf.paragraphs[0]
            p.text = feature
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['text_white']
            y_pos += 0.33
    
    # Integration with Nova
    nova_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    tf = nova_box.text_frame
    p = tf.paragraphs[0]
    p.text = "⚡ INTEGRATION: All agents connect to Nova AI & OneDevelopment systems"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_green']
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "PropertyFinder API  •  Bayut API  •  Internal CRM  •  Document Management"
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLORS['text_gray']
    p2.alignment = PP_ALIGN.CENTER

def add_week1_slide(prs):
    """Slide 5: Week 1 - Foundation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📅 WEEK 1: FOUNDATION (Days 1-7)"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_purple']
    
    # Target
    target_box = slide.shapes.add_textbox(Inches(7.2), Inches(0.3), Inches(2.5), Inches(0.5))
    tf = target_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯 5 Agents"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_green']
    
    agents = [
        ("Agent 1", "Nova Enhancement Agent", "Enhances Nova AI with deeper property knowledge", COLORS['accent_gold']),
        ("Agent 2", "Property Data Collector", "Auto-scrapes listings from PropertyFinder & Bayut", COLORS['accent_cyan']),
        ("Agent 3", "Client Inquiry Router", "Routes & prioritizes client inquiries to right teams", COLORS['accent_green']),
        ("Agent 4", "Document Summarizer", "Summarizes contracts, agreements, legal docs", COLORS['accent_orange']),
        ("Agent 5", "Meeting Notes Agent", "Auto-generates meeting notes & action items", COLORS['accent_pink']),
    ]
    
    y_pos = 1
    for agent, name, desc, color in agents:
        # Agent number
        num_box = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos), Inches(1.2), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = agent
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        # Name
        name_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos), Inches(2.8), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_white']
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos), Inches(5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['text_gray']
        
        y_pos += 0.55
    
    # Focus note
    focus_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = focus_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🔧 Week 1 Focus: Core infrastructure & customer-facing agents"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['accent_cyan']
    p.alignment = PP_ALIGN.CENTER

def add_week2_slide(prs):
    """Slide 6: Week 2 - Sales & Marketing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📅 WEEK 2: SALES & MARKETING (Days 8-14)"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_gold']
    
    # Target
    target_box = slide.shapes.add_textbox(Inches(7.2), Inches(0.3), Inches(2.5), Inches(0.5))
    tf = target_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯 5 Agents"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_green']
    
    agents = [
        ("Agent 6", "Lead Qualification Agent", "Scores & qualifies inbound property leads", COLORS['accent_gold']),
        ("Agent 7", "Market Analysis Agent", "UAE real estate trends & competitor analysis", COLORS['accent_cyan']),
        ("Agent 8", "Email Campaign Agent", "Personalizes property marketing emails", COLORS['accent_green']),
        ("Agent 9", "Social Media Agent", "Creates property posts for Instagram/LinkedIn", COLORS['accent_orange']),
        ("Agent 10", "Client Follow-up Agent", "Automated follow-ups with interested buyers", COLORS['accent_pink']),
    ]
    
    y_pos = 1
    for agent, name, desc, color in agents:
        num_box = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos), Inches(1.2), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = agent
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        name_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos), Inches(2.8), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_white']
        
        desc_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos), Inches(5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['text_gray']
        
        y_pos += 0.55
    
    focus_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = focus_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📈 Week 2 Focus: Revenue-generating agents for sales team"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['accent_gold']
    p.alignment = PP_ALIGN.CENTER

def add_week3_slide(prs):
    """Slide 7: Week 3 - Operations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📅 WEEK 3: OPERATIONS (Days 15-21)"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_green']
    
    # Target
    target_box = slide.shapes.add_textbox(Inches(7.2), Inches(0.3), Inches(2.5), Inches(0.5))
    tf = target_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯 5 Agents"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_green']
    
    agents = [
        ("Agent 11", "Project Tracker Agent", "Tracks construction milestones & deadlines", COLORS['accent_gold']),
        ("Agent 12", "Vendor Management Agent", "Manages supplier communications & orders", COLORS['accent_cyan']),
        ("Agent 13", "Quality Checklist Agent", "Construction quality inspection tracking", COLORS['accent_green']),
        ("Agent 14", "Budget Monitor Agent", "Tracks project expenses & budget alerts", COLORS['accent_orange']),
        ("Agent 15", "Compliance Agent", "UAE RERA regulations & permit tracking", COLORS['accent_pink']),
    ]
    
    y_pos = 1
    for agent, name, desc, color in agents:
        num_box = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos), Inches(1.2), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = agent
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        name_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos), Inches(2.8), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_white']
        
        desc_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos), Inches(5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['text_gray']
        
        y_pos += 0.55
    
    focus_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = focus_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🏗️ Week 3 Focus: Development & construction efficiency"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['accent_green']
    p.alignment = PP_ALIGN.CENTER

def add_week4_slide(prs):
    """Slide 8: Week 4 - HR & Internal"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📅 WEEK 4: HR & DEPLOYMENT (Days 22-30)"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_orange']
    
    # Target
    target_box = slide.shapes.add_textbox(Inches(7.2), Inches(0.3), Inches(2.5), Inches(0.5))
    tf = target_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎯 5 Agents"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_green']
    
    agents = [
        ("Agent 16", "HR Onboarding Agent", "New employee onboarding & training", COLORS['accent_gold']),
        ("Agent 17", "Knowledge Base Agent", "Internal docs search & policy lookup", COLORS['accent_cyan']),
        ("Agent 18", "IT Support Agent", "First-line IT troubleshooting", COLORS['accent_green']),
        ("Agent 19", "Performance Report Agent", "Weekly/monthly KPI reports", COLORS['accent_orange']),
        ("Agent 20", "Agent Orchestrator", "Meta-agent coordinating all 19 agents", COLORS['accent_pink']),
    ]
    
    y_pos = 1
    for agent, name, desc, color in agents:
        num_box = slide.shapes.add_textbox(Inches(0.4), Inches(y_pos), Inches(1.2), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = agent
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color
        
        name_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos), Inches(2.8), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['text_white']
        
        desc_box = slide.shapes.add_textbox(Inches(4.5), Inches(y_pos), Inches(5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['text_gray']
        
        y_pos += 0.55
    
    # Deployment note
    deploy_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(9), Inches(0.9))
    tf = deploy_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🚀 Week 4 Includes: Production deployment, testing, training"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['accent_orange']
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "All 20 agents integrated with OneDevelopment infrastructure"
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLORS['text_gray']
    p2.alignment = PP_ALIGN.CENTER

def add_agent_categories_slide(prs):
    """Slide 9: Agent Categories by Department"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🏢 AGENTS BY DEPARTMENT"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    
    categories = [
        {
            'name': '💼 SALES',
            'color': COLORS['accent_gold'],
            'agents': ['Lead Qualification', 'Client Follow-up', 'Client Inquiry Router']
        },
        {
            'name': '📣 MARKETING',
            'color': COLORS['accent_cyan'],
            'agents': ['Email Campaign', 'Social Media', 'Market Analysis']
        },
        {
            'name': '🏗️ OPERATIONS',
            'color': COLORS['accent_green'],
            'agents': ['Project Tracker', 'Vendor Mgmt', 'Quality Checklist', 'Budget Monitor']
        },
        {
            'name': '⚖️ LEGAL/COMPLIANCE',
            'color': COLORS['accent_orange'],
            'agents': ['Document Summarizer', 'Compliance Agent']
        },
        {
            'name': '👥 HR/ADMIN',
            'color': COLORS['accent_pink'],
            'agents': ['HR Onboarding', 'Meeting Notes', 'IT Support']
        },
        {
            'name': '📊 INTELLIGENCE',
            'color': COLORS['accent_purple'],
            'agents': ['Nova Enhancement', 'Property Data', 'Knowledge Base', 'Performance Report', 'Orchestrator']
        }
    ]
    
    positions = [
        (0.3, 1), (3.4, 1), (6.5, 1),
        (0.3, 2.7), (3.4, 2.7), (6.5, 2.7)
    ]
    
    for i, cat in enumerate(categories):
        x, y = positions[i]
        
        name_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = cat['name']
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = cat['color']
        
        agents_y = y + 0.35
        for agent in cat['agents']:
            agent_box = slide.shapes.add_textbox(Inches(x), Inches(agents_y), Inches(2.8), Inches(0.28))
            tf = agent_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {agent}"
            p.font.size = Pt(9)
            p.font.color.rgb = COLORS['text_white']
            agents_y += 0.26

def add_employee_benefits_slide(prs):
    """Slide 10: Benefits for Employees"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "👥 BENEFITS FOR EMPLOYEES"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    
    benefits = [
        ("🕐", "Save 2-3 hours daily", "Automate repetitive research & documentation tasks", COLORS['accent_gold']),
        ("📊", "Better decisions", "Real-time market data & competitor insights", COLORS['accent_cyan']),
        ("🤝", "Focus on relationships", "Let agents handle admin while you build client trust", COLORS['accent_green']),
        ("📚", "Instant knowledge", "Search internal policies, contracts, procedures in seconds", COLORS['accent_orange']),
        ("🎯", "Higher conversions", "Qualified leads with AI scoring & smart follow-ups", COLORS['accent_pink']),
    ]
    
    y_pos = 1.1
    for icon, title, desc, color in benefits:
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos), Inches(0.6), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(20)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1.3), Inches(y_pos), Inches(3), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(4.3), Inches(y_pos + 0.05), Inches(5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['text_gray']
        
        y_pos += 0.65

def add_roi_slide(prs):
    """Slide 11: ROI & Impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 EXPECTED IMPACT"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    
    metrics = [
        ("70%", "Faster", "lead response", COLORS['accent_gold']),
        ("500+", "Hours", "saved monthly", COLORS['accent_cyan']),
        ("40%", "More", "qualified leads", COLORS['accent_green']),
        ("24/7", "Support", "availability", COLORS['accent_orange']),
    ]
    
    x_positions = [0.4, 2.6, 4.8, 7]
    
    for i, (num, label1, label2, color) in enumerate(metrics):
        x = x_positions[i]
        
        num_box = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(2), Inches(0.8))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        
        label_box = slide.shapes.add_textbox(Inches(x), Inches(2), Inches(2), Inches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{label1}\n{label2}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['text_white']
        p.alignment = PP_ALIGN.CENTER
    
    # Department impacts
    impacts_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(0.4))
    tf = impacts_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Impact by Department:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_purple']
    
    impacts = [
        ("Sales:", "30% faster deal cycles with automated follow-ups"),
        ("Marketing:", "4x content output with AI-generated campaigns"),
        ("Operations:", "Real-time project visibility, 50% fewer delays"),
        ("HR:", "90% faster onboarding, instant policy answers"),
    ]
    
    y_pos = 3.35
    for dept, impact in impacts:
        dept_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(1.5), Inches(0.35))
        tf = dept_box.text_frame
        p = tf.paragraphs[0]
        p.text = dept
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent_gold']
        
        impact_box = slide.shapes.add_textbox(Inches(2.3), Inches(y_pos), Inches(7), Inches(0.35))
        tf = impact_box.text_frame
        p = tf.paragraphs[0]
        p.text = impact
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['text_white']
        
        y_pos += 0.4

def add_conclusion_slide(prs):
    """Slide 12: Conclusion & Next Steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🚀 LET'S TRANSFORM ONEDEVELOPMENT"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.alignment = PP_ALIGN.CENTER
    
    # Summary
    summary_box = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(8), Inches(0.8))
    tf = summary_box.text_frame
    p = tf.paragraphs[0]
    p.text = "20 AI agents tailored for OneDevelopment employees"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['accent_gold']
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Built with DeepAgents + LangSmith + Cursor in just 30 days"
    p2.font.size = Pt(14)
    p2.font.color.rgb = COLORS['text_gray']
    p2.alignment = PP_ALIGN.CENTER
    
    # Next steps
    next_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(7), Inches(2))
    tf = next_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📋 Next Steps:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_cyan']
    
    steps = [
        "1. Approve the 30-day project timeline",
        "2. Identify department champions for each agent",
        "3. Finalize API access & system integrations",
        "4. Kick off Week 1: Foundation phase",
        "5. Weekly demos & feedback sessions"
    ]
    
    for step in steps:
        p2 = tf.add_paragraph()
        p2.text = f"    {step}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLORS['text_white']
    
    # Call to action
    cta_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.5))
    tf = cta_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ready to empower every OneDevelopment employee with AI?"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent_purple']
    p.alignment = PP_ALIGN.CENTER

def create_presentation():
    """Create the full PowerPoint presentation"""
    prs = Presentation()
    
    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Add all slides
    add_title_slide(prs)
    add_agenda_slide(prs)
    add_why_agents_slide(prs)
    add_tech_stack_slide(prs)
    add_week1_slide(prs)
    add_week2_slide(prs)
    add_week3_slide(prs)
    add_week4_slide(prs)
    add_agent_categories_slide(prs)
    add_employee_benefits_slide(prs)
    add_roi_slide(prs)
    add_conclusion_slide(prs)
    
    # Save presentation
    output_path = "/home/ec2-user/OneDevelopment-Agent/20_Agents_30_Days_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created successfully!")
    print(f"📁 Saved to: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    create_presentation()
