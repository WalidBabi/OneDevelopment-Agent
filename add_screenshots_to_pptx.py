#!/usr/bin/env python3
"""
Add Nova AI screenshots to the PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Color palette - OneDevelopment Purple Theme
COLORS = {
    'dark_bg': RGBColor(26, 13, 48),
    'accent_purple': RGBColor(150, 107, 252),
    'accent_gold': RGBColor(255, 193, 7),
    'accent_cyan': RGBColor(34, 211, 238),
    'text_white': RGBColor(255, 255, 255),
    'text_gray': RGBColor(180, 180, 195),
}

SCREENSHOT_DIR = "/home/ec2-user/OneDevelopment-Agent/screenshots"
PPTX_PATH = "/home/ec2-user/OneDevelopment-Agent/20_Agents_30_Days_Presentation.pptx"

def set_slide_background(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_nova_demo_slide(prs, screenshot_path, title, subtitle, position="center"):
    """Add a slide with a Nova screenshot"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.6), Inches(9), Inches(0.3))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['text_gray']
    
    # Add screenshot image
    if os.path.exists(screenshot_path):
        if position == "center":
            # Centered image
            slide.shapes.add_picture(
                screenshot_path,
                Inches(0.8),
                Inches(1),
                width=Inches(8.4),
                height=Inches(4.3)
            )
        elif position == "mobile":
            # Mobile view - smaller, centered
            slide.shapes.add_picture(
                screenshot_path,
                Inches(3.5),
                Inches(0.9),
                width=Inches(3),
                height=Inches(4.5)
            )
        print(f"  ✓ Added: {os.path.basename(screenshot_path)}")
    else:
        print(f"  ✗ Not found: {screenshot_path}")
    
    return slide

def add_intro_nova_slide(prs):
    """Add an introduction slide for Nova section"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS['dark_bg'])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🤖 MEET NOVA"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_white']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "OneDevelopment's AI Assistant - Already Live!"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['accent_gold']
    p.alignment = PP_ALIGN.CENTER
    
    # URL
    url_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.4))
    tf = url_box.text_frame
    p = tf.paragraphs[0]
    p.text = "<YOUR_SERVER_URL>:3000/"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['accent_cyan']
    p.alignment = PP_ALIGN.CENTER
    
    # Features
    features_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(7), Inches(1))
    tf = features_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GPT-4 Powered  •  Property Knowledge  •  24/7 Available  •  Multi-Source Intelligence"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['text_gray']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    """Add screenshots to the existing presentation"""
    
    print("📊 Opening presentation...")
    prs = Presentation(PPTX_PATH)
    
    # Get current slide count
    original_count = len(prs.slides)
    print(f"   Current slides: {original_count}")
    
    print("\n📸 Adding Nova screenshots...")
    
    # Add Nova intro slide (after conclusion, before we add screenshots)
    add_intro_nova_slide(prs)
    print("  ✓ Added: Nova intro slide")
    
    # Add screenshot slides
    screenshots = [
        ("nova_landing.png", "🖥️ NOVA - LANDING PAGE", "Clean, modern interface welcoming OneDevelopment visitors"),
        ("nova_chat_typing.png", "💬 NOVA - CHAT INTERFACE", "Natural language conversation about properties"),
        ("nova_chat_response.png", "🧠 NOVA - AI RESPONSE", "Intelligent, context-aware property information"),
        ("nova_mobile.png", "📱 NOVA - MOBILE VIEW", "Responsive design for on-the-go access"),
    ]
    
    for filename, title, subtitle in screenshots:
        path = os.path.join(SCREENSHOT_DIR, filename)
        position = "mobile" if "mobile" in filename else "center"
        add_nova_demo_slide(prs, path, title, subtitle, position)
    
    # Save updated presentation
    prs.save(PPTX_PATH)
    
    final_count = len(prs.slides)
    print(f"\n✅ Presentation updated!")
    print(f"   Original slides: {original_count}")
    print(f"   Added slides: {final_count - original_count}")
    print(f"   Total slides: {final_count}")
    print(f"📁 Saved to: {PPTX_PATH}")

if __name__ == "__main__":
    main()

